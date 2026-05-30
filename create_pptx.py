from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

BLUE = RGBColor(0x1A, 0x3C, 0x6E)
DARK = RGBColor(0x2D, 0x2D, 0x2D)
GRAY = RGBColor(0x66, 0x66, 0x66)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT = RGBColor(0x4A, 0x6C, 0xF7)
GREEN = RGBColor(0x28, 0xA7, 0x45)
RED = RGBColor(0xDC, 0x35, 0x45)
LIGHT_BG = RGBColor(0xF0, 0xF2, 0xF5)
CARD_BG = RGBColor(0xFF, 0xFF, 0xFF)

def add_bg(slide, color=LIGHT_BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape(slide, left, top, width, height, fill_color, line_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=18, color=DARK, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_bullet_text(slide, left, top, width, height, items, font_size=16, color=DARK, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = font_name
        p.space_after = Pt(6)
        p.level = 0
    return txBox

# ============================================================
# SLIDE 1: TITLE
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BLUE)

add_text_box(slide, Inches(1), Inches(1.5), Inches(11), Inches(1.5),
    "Deepfake Video Detection", font_size=48, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1), Inches(3.2), Inches(11), Inches(1),
    "Using ResNet50 Feature Extraction & XGBoost Classification", font_size=24, color=RGBColor(0xBB, 0xCC, 0xFF), alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1), Inches(5), Inches(11), Inches(0.8),
    "FaceForensics++  |  OpenCV DNN Face Detector  |  Flask Web App", font_size=18, color=RGBColor(0x88, 0xAA, 0xDD), alignment=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 2: PROBLEM STATEMENT
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(5), Inches(0.8),
    "Problem Statement", font_size=36, color=BLUE, bold=True)

# Card 1
card = add_shape(slide, Inches(0.8), Inches(1.4), Inches(5.5), Inches(2.4), CARD_BG, RGBColor(0xDD, 0xDD, 0xDD))
add_text_box(slide, Inches(1.1), Inches(1.6), Inches(5), Inches(0.5),
    "The Threat of Deepfakes", font_size=22, color=DARK, bold=True)
add_bullet_text(slide, Inches(1.1), Inches(2.2), Inches(5), Inches(1.4), [
    "AI-generated videos (deepfakes) are increasingly realistic",
    "Used for misinformation, fraud, and identity theft",
    "Manual detection is impossible at scale",
    "Need automated, accurate detection systems",
], font_size=16, color=GRAY)

# Card 2
card = add_shape(slide, Inches(7), Inches(1.4), Inches(5.5), Inches(2.4), CARD_BG, RGBColor(0xDD, 0xDD, 0xDD))
add_text_box(slide, Inches(7.3), Inches(1.6), Inches(5), Inches(0.5),
    "Project Objective", font_size=22, color=DARK, bold=True)
add_bullet_text(slide, Inches(7.3), Inches(2.2), Inches(5), Inches(1.4), [
    "Build a deepfake video detector using deep learning",
    "Classify videos as REAL or FAKE with high accuracy",
    "Deploy as a user-friendly web application (Flask)",
    "Handle the FaceForensics++ benchmark dataset",
], font_size=16, color=GRAY)

# Bottom insight
card = add_shape(slide, Inches(0.8), Inches(4.2), Inches(11.7), Inches(2.5), CARD_BG, RGBColor(0xDD, 0xDD, 0xDD))
add_text_box(slide, Inches(1.1), Inches(4.4), Inches(11), Inches(0.5),
    "Key Challenge: Avoiding Detection Artifacts", font_size=22, color=DARK, bold=True)
add_bullet_text(slide, Inches(1.1), Inches(5.0), Inches(11), Inches(1.5), [
    "Using the SAME face detector for both real and fake extraction during training",
    "Prevents the model from learning extraction differences instead of real vs. fake features",
    "Previous models achieved 99% accuracy by cheating (learning which detector was used, not actual deepfake artifacts)",
    "Our approach uses OpenCV DNN for consistent extraction across both classes",
], font_size=16, color=GRAY)

# ============================================================
# SLIDE 3: DATASET
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(5), Inches(0.8),
    "Dataset: FaceForensics++", font_size=36, color=BLUE, bold=True)

# Stats cards
for i, (title, val, desc) in enumerate([
    ("Total Videos", "400", "200 Real + 200 Fake"),
    ("Real Videos", "200", "Original YouTube interviews"),
    ("Fake Videos", "200", "FaceSwap generated deepfakes"),
    ("Train/Val/Test", "274 / 40 / 79", "Video-level split (no leakage)"),
]):
    left = Inches(0.8 + i * 3.1)
    card = add_shape(slide, left, Inches(1.5), Inches(2.8), Inches(2.2), CARD_BG, RGBColor(0xDD, 0xDD, 0xDD))
    add_text_box(slide, left + Inches(0.2), Inches(1.7), Inches(2.4), Inches(0.5),
        title, font_size=16, color=GRAY)
    add_text_box(slide, left + Inches(0.2), Inches(2.2), Inches(2.4), Inches(0.7),
        val, font_size=36, color=BLUE, bold=True)
    add_text_box(slide, left + Inches(0.2), Inches(3.0), Inches(2.4), Inches(0.5),
        desc, font_size=14, color=GRAY)

# Bottom info
add_shape(slide, Inches(0.8), Inches(4.2), Inches(11.7), Inches(2.5), CARD_BG, RGBColor(0xDD, 0xDD, 0xDD))
add_text_box(slide, Inches(1.1), Inches(4.4), Inches(11), Inches(0.5),
    "Pre-processing Pipeline", font_size=22, color=DARK, bold=True)
add_bullet_text(slide, Inches(1.1), Inches(5.0), Inches(11), Inches(1.5), [
    "Each video: ~300 frames, ~5 seconds duration, 5.1 MB average size",
    "15 frames sampled per video using uniform spacing",
    "OpenCV DNN face detection on each sampled frame",
    "Extracted faces resized to 224x224 for ResNet50 input",
    "Features extracted from ResNet50 (no training, just forward pass)",
], font_size=16, color=GRAY)

# ============================================================
# SLIDE 4: SYSTEM ARCHITECTURE
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(5), Inches(0.8),
    "System Architecture", font_size=36, color=BLUE, bold=True)

# Pipeline boxes
stages = [
    ("Input\nVideo", RGBColor(0x4A, 0x6C, 0xF7)),
    ("Face Detection\n(OpenCV DNN)", RGBColor(0x6C, 0x8C, 0xF7)),
    ("Feature Extraction\n(ResNet50 CNN)", RGBColor(0x8C, 0xAC, 0xF7)),
    ("Classification\n(XGBoost)", RGBColor(0x28, 0xA7, 0x45)),
    ("Prediction\n(REAL / FAKE)", RGBColor(0xDC, 0x35, 0x45)),
]

for i, (label, color) in enumerate(stages):
    left = Inches(0.8 + i * 2.5)
    card = add_shape(slide, left, Inches(1.8), Inches(2.2), Inches(2.5), color)
    add_text_box(slide, left + Inches(0.1), Inches(2.2), Inches(2.0), Inches(1.8),
        label, font_size=18, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    # Arrow between boxes
    if i < len(stages) - 1:
        add_text_box(slide, left + Inches(2.1), Inches(2.7), Inches(0.5), Inches(0.5),
            ">>", font_size=24, color=GRAY, bold=True, alignment=PP_ALIGN.CENTER)

# Bottom detail cards
details = [
    ("OpenCV DNN (ResNet-10 SSD)", "Face Detector", [
        "Pre-trained Caffe model",
        "Confidence threshold: 0.5",
        "300x300 input blob size",
        "Detects best face per frame",
    ]),
    ("ResNet50 (ImageNet)", "Feature Extractor", [
        "Pretrained on ImageNet-1K",
        "Global average pooling",
        "Output: 2048-dim feature vector",
        "Frozen weights (no fine-tuning)",
    ]),
    ("XGBoost Classifier", "Classifier", [
        "300 estimators, max depth 6",
        "Learning rate: 0.05",
        "L1 (alpha=1.0) + L2 (lambda=2.0) reg",
        "Subsample 0.8, colsample 0.8",
    ]),
]

for i, (title, subtitle, items) in enumerate(details):
    left = Inches(0.8 + i * 4.1)
    card = add_shape(slide, left, Inches(4.8), Inches(3.8), Inches(2.5), CARD_BG, RGBColor(0xDD, 0xDD, 0xDD))
    add_text_box(slide, left + Inches(0.3), Inches(5.0), Inches(3.3), Inches(0.4),
        title, font_size=18, color=BLUE, bold=True)
    add_text_box(slide, left + Inches(0.3), Inches(5.4), Inches(3.3), Inches(0.3),
        subtitle, font_size=13, color=GRAY)
    add_bullet_text(slide, left + Inches(0.3), Inches(5.8), Inches(3.3), Inches(1.3),
        items, font_size=14, color=DARK)

# ============================================================
# SLIDE 5: FACE DETECTION DETAIL
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(5), Inches(0.8),
    "Face Detection: OpenCV DNN", font_size=36, color=BLUE, bold=True)

# Left card
card = add_shape(slide, Inches(0.8), Inches(1.5), Inches(5.8), Inches(5.3), CARD_BG, RGBColor(0xDD, 0xDD, 0xDD))
add_text_box(slide, Inches(1.1), Inches(1.7), Inches(5.3), Inches(0.5),
    "ResNet-10 SSD Face Detector", font_size=22, color=DARK, bold=True)
add_bullet_text(slide, Inches(1.1), Inches(2.3), Inches(5.3), Inches(2.0), [
    "Single Shot Multibox Detector (SSD) framework",
    "ResNet-10 base architecture (lightweight)",
    "Pre-trained on WIDER Face dataset",
    "Outputs bounding boxes with confidence scores",
    "Picks highest-confidence face per frame",
], font_size=16, color=GRAY)

add_text_box(slide, Inches(1.1), Inches(4.3), Inches(5.3), Inches(0.5),
    "Parameters", font_size=22, color=DARK, bold=True)
add_bullet_text(slide, Inches(1.1), Inches(4.9), Inches(5.3), Inches(1.5), [
    "Input blob: 300x300 pixels",
    "Mean subtraction: (104, 117, 123)",
    "Scale factor: 1.0",
    "Confidence threshold: 0.5 (training), 0.2 (inference)",
    "Detects faces in ~0.1s per frame on GPU",
], font_size=16, color=GRAY)

# Right card
card = add_shape(slide, Inches(7), Inches(1.5), Inches(5.5), Inches(5.3), CARD_BG, RGBColor(0xDD, 0xDD, 0xDD))
add_text_box(slide, Inches(7.3), Inches(1.7), Inches(5), Inches(0.5),
    "Why OpenCV DNN?", font_size=22, color=DARK, bold=True)
add_bullet_text(slide, Inches(7.3), Inches(2.3), Inches(5), Inches(2.0), [
    "Consistent: Same detector for BOTH fake and real videos",
    "Fast: GPU-accelerated via CUDA backend",
    "Lightweight: Single .caffemodel file (10.6 MB)",
    "No extra dependencies (OpenCV built-in)",
    "Works well on 150x150+ faces in interview videos",
], font_size=16, color=GRAY)

add_text_box(slide, Inches(7.3), Inches(4.3), Inches(5), Inches(0.5),
    "Key Design Decision", font_size=22, color=RED, bold=True)
add_text_box(slide, Inches(7.3), Inches(4.9), Inches(5), Inches(1.8),
    "Using the same face detector for both classes during training is critical. Prior attempts used MTCNN for fake and OpenCV DNN for real, achieving 99% accuracy. This was ARTIFACT LEARNING - the model learned which detector was used, not actual deepfake features. Our approach produces a genuine, generalizable model.",
    font_size=15, color=GRAY)

# ============================================================
# SLIDE 6: FEATURE EXTRACTION (ResNet50)
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(5), Inches(0.8),
    "Feature Extraction: ResNet50", font_size=36, color=BLUE, bold=True)

card = add_shape(slide, Inches(0.8), Inches(1.5), Inches(11.7), Inches(2.5), CARD_BG, RGBColor(0xDD, 0xDD, 0xDD))
add_text_box(slide, Inches(1.1), Inches(1.7), Inches(11), Inches(0.5),
    "ResNet50 Architecture Overview", font_size=22, color=DARK, bold=True)
add_bullet_text(slide, Inches(1.1), Inches(2.3), Inches(11), Inches(1.5), [
    "50-layer deep residual network with skip connections (bottleneck blocks)",
    "Input: 224x224 RGB face image | Pre-processing: Normalize to mean=(0.485,0.456,0.406), std=(0.229,0.224,0.225)",
    "Last fully-connected layer removed -> Global Average Pooling -> 2048-dimensional feature vector",
    "Pre-trained on ImageNet (1.2M images, 1000 classes) - used as FIXED feature extractor (no fine-tuning)",
    "CUDA-accelerated: processes batch of 15 frames in ~0.3 seconds",
], font_size=16, color=GRAY)

card = add_shape(slide, Inches(0.8), Inches(4.3), Inches(5.8), Inches(2.5), CARD_BG, RGBColor(0xDD, 0xDD, 0xDD))
add_text_box(slide, Inches(1.1), Inches(4.5), Inches(5.3), Inches(0.5),
    "Why ResNet50?", font_size=22, color=DARK, bold=True)
add_bullet_text(slide, Inches(1.1), Inches(5.1), Inches(5.3), Inches(1.5), [
    "Proven feature extraction quality (ImageNet SOTA)", 
    "2048-dim vectors capture high-level facial features",
    "Skip connections prevent vanishing gradients",
    "Well-understood model with PyTorch support",
], font_size=16, color=GRAY)

card = add_shape(slide, Inches(7), Inches(4.3), Inches(5.5), Inches(2.5), CARD_BG, RGBColor(0xDD, 0xDD, 0xDD))
add_text_box(slide, Inches(7.3), Inches(4.5), Inches(5), Inches(0.5),
    "Feature Processing Pipeline", font_size=22, color=DARK, bold=True)
add_bullet_text(slide, Inches(7.3), Inches(5.1), Inches(5), Inches(1.5), [
    "Face image -> resize to 224x224 -> normalize",
    "Batch all 15 frames -> forward through ResNet50",
    "Average Pool (7x7 -> 1x1) -> flatten to 2048-dim",
    "Per-video: 15 x 2048 feature matrix -> average -> XGBoost",
], font_size=16, color=GRAY)

# ============================================================
# SLIDE 7: XGBOOST CLASSIFIER
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(5), Inches(0.8),
    "Classifier: XGBoost", font_size=36, color=BLUE, bold=True)

card = add_shape(slide, Inches(0.8), Inches(1.5), Inches(5.8), Inches(5.3), CARD_BG, RGBColor(0xDD, 0xDD, 0xDD))
add_text_box(slide, Inches(1.1), Inches(1.7), Inches(5.3), Inches(0.5),
    "XGBoost Hyperparameters", font_size=22, color=DARK, bold=True)

params = [
    ("n_estimators", "300", "Number of boosting rounds"),
    ("max_depth", "6", "Tree depth (controls overfitting)"),
    ("learning_rate", "0.05", "Shrinkage step size"),
    ("subsample", "0.8", "Row sampling per tree"),
    ("colsample_bytree", "0.8", "Column sampling per tree"),
    ("reg_lambda", "2.0", "L2 regularization weight"),
    ("reg_alpha", "1.0", "L1 regularization weight"),
    ("eval_metric", "logloss", "Binary classification loss"),
]

for i, (param, value, desc) in enumerate(params):
    y = Inches(2.3) + Inches(i * 0.55)
    add_text_box(slide, Inches(1.1), y, Inches(1.8), Inches(0.4),
        param, font_size=14, color=BLUE, bold=True)
    add_text_box(slide, Inches(3.0), y, Inches(1.2), Inches(0.4),
        value, font_size=14, color=DARK, bold=True)
    add_text_box(slide, Inches(4.2), y, Inches(2.2), Inches(0.4),
        desc, font_size=13, color=GRAY)

card = add_shape(slide, Inches(7), Inches(1.5), Inches(5.5), Inches(5.3), CARD_BG, RGBColor(0xDD, 0xDD, 0xDD))
add_text_box(slide, Inches(7.3), Inches(1.7), Inches(5), Inches(0.5),
    "Why XGBoost over Deep Classifier?", font_size=22, color=DARK, bold=True)
add_bullet_text(slide, Inches(7.3), Inches(2.3), Inches(5), Inches(2.0), [
    "Works well with limited data (400 videos, 4000 frames)",
    "Handles 2048-dim feature vectors efficiently",
    "Built-in regularization prevents overfitting",
    "Faster training than fine-tuning a full CNN",
    "Tree-based ensembles are interpretable",
    "Excellent for tabular data (feature vectors)",
], font_size=16, color=GRAY)

add_text_box(slide, Inches(7.3), Inches(4.5), Inches(5), Inches(0.5),
    "Training Configuration", font_size=22, color=DARK, bold=True)
add_bullet_text(slide, Inches(7.3), Inches(5.1), Inches(5), Inches(1.5), [
    "Early stopping rounds from validation",
    "Video-level train/val/test split (no leakage)",
    "Frame-level training, video-level evaluation",
    "Average frame probabilities for video prediction",
], font_size=16, color=GRAY)

# ============================================================
# SLIDE 8: RESULTS - METRICS TABLE
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(5), Inches(0.8),
    "Training Results", font_size=36, color=BLUE, bold=True)

# Metrics table
card = add_shape(slide, Inches(0.8), Inches(1.5), Inches(11.7), Inches(2.5), CARD_BG, RGBColor(0xDD, 0xDD, 0xDD))

# Table header
headers = ["Metric", "Train (Frame)", "Validation (Frame)", "Test (Video-Level)"]
for i, h in enumerate(headers):
    left = Inches(1.0 + i * 2.9)
    add_text_box(slide, left, Inches(1.7), Inches(2.6), Inches(0.5),
        h, font_size=16, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    hdr_bg = add_shape(slide, left - Inches(0.1), Inches(1.65), Inches(2.7), Inches(0.5), BLUE)

rows = [
    ("Accuracy", "1.0000", "0.7177", "0.7848"),
    ("AUC-ROC", "1.0000", "0.7791", "0.8622"),
    ("F1-Score", "1.0000", "0.7533", "0.7848"),
]

for r, (label, train, val, test) in enumerate(rows):
    y = Inches(2.3) + Inches(r * 0.6)
    bg_color = RGBColor(0xF8, 0xF9, 0xFA) if r % 2 == 0 else CARD_BG
    add_shape(slide, Inches(0.9), y, Inches(2.6), Inches(0.5), bg_color)
    add_text_box(slide, Inches(1.1), y, Inches(2.3), Inches(0.5),
        label, font_size=16, color=DARK, bold=True)
    
    for c, val_text in enumerate([train, val, test]):
        left = Inches(3.5 + c * 2.9)
        add_shape(slide, left, y, Inches(2.6), Inches(0.5), bg_color)
        color = GREEN if r == 0 and c < 2 else (GREEN if r == 1 and c == 2 else DARK)
        add_text_box(slide, left + Inches(0.2), y, Inches(2.2), Inches(0.5),
            val_text, font_size=16, color=color, bold=True, alignment=PP_ALIGN.CENTER)

# Classification report
card = add_shape(slide, Inches(0.8), Inches(4.3), Inches(5.8), Inches(2.5), CARD_BG, RGBColor(0xDD, 0xDD, 0xDD))
add_text_box(slide, Inches(1.1), Inches(4.5), Inches(5.3), Inches(0.5),
    "Video-Level Classification Report", font_size=22, color=DARK, bold=True)

# Sub-table
sub_headers = ["Class", "Precision", "Recall", "F1-Score", "Support"]
for i, h in enumerate(sub_headers):
    left = Inches(1.1 + i * 1.1)
    add_text_box(slide, left, Inches(5.1), Inches(1.0), Inches(0.4),
        h, font_size=13, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_shape(slide, left - Inches(0.05), Inches(5.05), Inches(1.05), Inches(0.35), BLUE)

sub_rows = [
    ("Real", "0.78", "0.79", "0.78", "39"),
    ("Fake", "0.79", "0.78", "0.78", "40"),
]
for r, (cls, prec, rec, f1, supp) in enumerate(sub_rows):
    y = Inches(5.5) + Inches(r * 0.5)
    bg = RGBColor(0xF8, 0xF9, 0xFA) if r % 2 == 0 else CARD_BG
    vals = [cls, prec, rec, f1, supp]
    for i, v in enumerate(vals):
        left = Inches(1.1 + i * 1.1)
        add_shape(slide, left - Inches(0.05), y, Inches(1.05), Inches(0.4), bg)
        add_text_box(slide, left, y + Pt(2), Inches(0.9), Inches(0.35),
            v, font_size=14, color=DARK, bold=(i == 0), alignment=PP_ALIGN.CENTER)

# Insight card
card = add_shape(slide, Inches(7), Inches(4.3), Inches(5.5), Inches(2.5), CARD_BG, RGBColor(0xDD, 0xDD, 0xDD))
add_text_box(slide, Inches(7.3), Inches(4.5), Inches(5), Inches(0.5),
    "Key Observations", font_size=22, color=DARK, bold=True)
add_bullet_text(slide, Inches(7.3), Inches(5.1), Inches(5), Inches(1.5), [
    "Perfect training accuracy (1.0) indicates model capacity is sufficient",
    "Val-test gap (72% vs 78%) suggests some overfitting",
    "AUC of 0.8622 shows strong discriminative ability",
    "Balanced precision/recall across both classes",
    "78% accuracy is HONEST - no extraction artifact cheating",
], font_size=16, color=GRAY)

# ============================================================
# SLIDE 9: CONFUSION MATRIX
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(5), Inches(0.8),
    "Confusion Matrix (Video-Level Test)", font_size=36, color=BLUE, bold=True)

# Insert confusion matrix image
img_path = "/home/syed/BNU_PROJECT/models/confusion_matrix.png"
if os.path.exists(img_path):
    slide.shapes.add_picture(img_path, Inches(1.5), Inches(1.5), Inches(5.5), Inches(5.0))
else:
    add_text_box(slide, Inches(1.5), Inches(3), Inches(5), Inches(1),
        "Confusion Matrix image not found\nRun training to generate it",
        font_size=20, color=RED, alignment=PP_ALIGN.CENTER)

# Analysis side card
card = add_shape(slide, Inches(7.5), Inches(1.5), Inches(5.0), Inches(5.3), CARD_BG, RGBColor(0xDD, 0xDD, 0xDD))
add_text_box(slide, Inches(7.8), Inches(1.7), Inches(4.5), Inches(0.5),
    "Confusion Matrix Analysis", font_size=22, color=DARK, bold=True)

add_text_box(slide, Inches(7.8), Inches(2.4), Inches(4.5), Inches(0.4),
    "Correct Predictions", font_size=18, color=GREEN, bold=True)
add_bullet_text(slide, Inches(7.8), Inches(2.9), Inches(4.5), Inches(0.8), [
    "Real videos correctly classified: 31 (True Negatives)",
    "Fake videos correctly classified: 31 (True Positives)",
], font_size=16, color=DARK)

add_text_box(slide, Inches(7.8), Inches(3.8), Inches(4.5), Inches(0.4),
    "Misclassifications", font_size=18, color=RED, bold=True)
add_bullet_text(slide, Inches(7.8), Inches(4.3), Inches(4.5), Inches(0.8), [
    "Real videos misclassified as Fake: 8 (False Positives)",
    "Fake videos misclassified as Real: 9 (False Negatives)",
], font_size=16, color=DARK)

add_text_box(slide, Inches(7.8), Inches(5.2), Inches(4.5), Inches(0.4),
    "Overall", font_size=18, color=BLUE, bold=True)
add_text_box(slide, Inches(7.8), Inches(5.7), Inches(4.5), Inches(0.8),
    "62 correct out of 79 test videos (78.48%)\n17 misclassified (21.52%)",
    font_size=16, color=DARK)

# ============================================================
# SLIDE 10: TRAINING PROGRESS
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(5), Inches(0.8),
    "Training Progress & Model Convergence", font_size=36, color=BLUE, bold=True)

# Training log loss card
card = add_shape(slide, Inches(0.8), Inches(1.5), Inches(5.8), Inches(5.3), CARD_BG, RGBColor(0xDD, 0xDD, 0xDD))
add_text_box(slide, Inches(1.1), Inches(1.7), Inches(5.3), Inches(0.5),
    "XGBoost Training Log-Loss", font_size=22, color=DARK, bold=True)

logloss_data = [
    ("Round 0", "0.66877", "0.67926"),
    ("Round 50", "0.20737", "0.57595"),
    ("Round 100", "0.10049", "0.58798"),
    ("Round 150", "0.05803", "0.61461"),
    ("Round 200", "0.03874", "0.63489"),
    ("Round 250", "0.02842", "0.64510"),
    ("Round 299", "0.02235", "0.66242"),
]

# Mini table
sub_headers = ["Boosting Round", "Train LogLoss", "Val LogLoss"]
for i, h in enumerate(sub_headers):
    left = Inches(1.4 + i * 1.6)
    add_text_box(slide, left, Inches(2.3), Inches(1.5), Inches(0.4),
        h, font_size=13, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_shape(slide, left - Inches(0.05), Inches(2.25), Inches(1.55), Inches(0.35), BLUE)

for r, (round_name, train_loss, val_loss) in enumerate(logloss_data):
    y = Inches(2.7) + Inches(r * 0.55)
    bg = RGBColor(0xF8, 0xF9, 0xFA) if r % 2 == 0 else CARD_BG
    vals = [round_name, train_loss, val_loss]
    colors = [DARK, GREEN, RED]
    for i, v in enumerate(vals):
        left = Inches(1.4 + i * 1.6)
        add_shape(slide, left - Inches(0.05), y, Inches(1.55), Inches(0.45), bg)
        add_text_box(slide, left, y + Pt(2), Inches(1.4), Inches(0.4),
            v, font_size=14, color=colors[i], bold=True, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1.1), Inches(6.3), Inches(5.3), Inches(0.3),
    "Training log-loss decreases steadily, validation stabilizes around 0.58-0.66",
    font_size=13, color=GRAY)

# Stats card
card = add_shape(slide, Inches(7), Inches(1.5), Inches(5.5), Inches(5.3), CARD_BG, RGBColor(0xDD, 0xDD, 0xDD))
add_text_box(slide, Inches(7.3), Inches(1.7), Inches(5), Inches(0.5),
    "Dataset Split Statistics", font_size=22, color=DARK, bold=True)

stats_items = [
    "Total videos processed: 393 (2 failed face detection)",
    "Fake videos used: 198 of 200",
    "Real videos used: 195 of 200",
    "Training videos: 274 (3,196 frames)",
    "Validation videos: 40 (464 frames)",
    "Test videos: 79 (balanced: 39 real + 40 fake)",
    "",
    "Training time: ~21 minutes total",
    "  - Fake extraction: 10 min 5 sec",
    "  - Real extraction: 9 min 44 sec",
    "  - XGBoost training: ~52 sec",
    "",
    "Inference speed: ~3 sec/video",
]
add_bullet_text(slide, Inches(7.3), Inches(2.3), Inches(5), Inches(4.3),
    stats_items, font_size=15, color=DARK)

# ============================================================
# SLIDE 11: WEB APPLICATION
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(5), Inches(0.8),
    "Web Application: Flask Deployment", font_size=36, color=BLUE, bold=True)

card = add_shape(slide, Inches(0.8), Inches(1.5), Inches(5.8), Inches(5.3), CARD_BG, RGBColor(0xDD, 0xDD, 0xDD))
add_text_box(slide, Inches(1.1), Inches(1.7), Inches(5.3), Inches(0.5),
    "App Features", font_size=22, color=DARK, bold=True)
add_bullet_text(slide, Inches(1.1), Inches(2.3), Inches(5.3), Inches(4.0), [
    "Upload video files (mp4, avi, mov, mkv, webm, flv)",
    "Max upload size: 500 MB",
    "Automatic face detection using OpenCV DNN",
    "25 uniformly sampled frames per video",
    "Per-frame face detection results shown",
    "Final prediction: REAL or FAKE with confidence %",
    "Video metadata: duration, FPS, total frames",
    "Temporary file cleanup after processing",
], font_size=16, color=GRAY)

card = add_shape(slide, Inches(7), Inches(1.5), Inches(5.5), Inches(5.3), CARD_BG, RGBColor(0xDD, 0xDD, 0xDD))
add_text_box(slide, Inches(7.3), Inches(1.7), Inches(5), Inches(0.5),
    "Technical Stack", font_size=22, color=DARK, bold=True)
add_bullet_text(slide, Inches(7.3), Inches(2.3), Inches(5), Inches(4.0), [
    "Framework: Flask (Python Web Framework)",
    "Frontend: HTML5, CSS3 (responsive design)",
    "Face Detection: OpenCV DNN (ResNet-10 SSD)",
    "Feature Extraction: PyTorch ResNet50",
    "Classifier: XGBoost (joblib serialized)",
    "Models loaded at startup (cold start: ~5s)",
    "Running on http://127.0.0.1:5000",
    "Works with GPU (CUDA) or CPU fallback",
], font_size=16, color=GRAY)

# ============================================================
# SLIDE 12: END-TO-END TEST
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(5), Inches(0.8),
    "End-to-End Verification", font_size=36, color=BLUE, bold=True)

# Test results cards
test_results = [
    ("Fake Video #1", "01_02__outside_talking_still_laughing", "FAKE", "96.91%", True),
    ("Fake Video #2", "01_02__walk_down_hall_angry", "FAKE", "95.19%", True),
    ("Fake Video #3", "01_03__hugging_happy", "REAL", "29.90%", False),
    ("Fake Video #4", "01_03__podium_speech", "FAKE", "98.33%", True),
    ("Fake Video #5", "01_03__talking_against_wall", "FAKE", "98.25%", True),
    ("Real Video #1", "01__exit_phone_room", "REAL", "3.96%", True),
    ("Real Video #2", "01__hugging_happy", "REAL", "14.12%", True),
    ("Real Video #3", "01__kitchen_pan", "REAL", "1.75%", True),
    ("Real Video #4", "01__kitchen_still", "REAL", "1.33%", True),
    ("Real Video #5", "01__meeting_serious", "REAL", "33.75%", True),
]

for i, (title, name, pred, conf, correct) in enumerate(test_results):
    col = i % 2
    row = i // 2
    left = Inches(0.8 + col * 6.2)
    top = Inches(1.5 + row * 1.15)
    
    bg_color = RGBColor(0xD4, 0xED, 0xDA) if pred == "REAL" else RGBColor(0xF8, 0xD7, 0xDA)
    card = add_shape(slide, left, top, Inches(5.9), Inches(1.0), bg_color)
    
    add_text_box(slide, left + Inches(0.2), top + Inches(0.05), Inches(3.5), Inches(0.35),
        title, font_size=14, color=DARK, bold=True)
    add_text_box(slide, left + Inches(0.2), top + Inches(0.4), Inches(4.0), Inches(0.3),
        name, font_size=11, color=GRAY)
    add_text_box(slide, left + Inches(4.0), top + Inches(0.1), Inches(1.8), Inches(0.4),
        pred, font_size=18, color=GREEN if pred == "REAL" else RED, bold=True, alignment=PP_ALIGN.RIGHT)
    add_text_box(slide, left + Inches(4.0), top + Inches(0.55), Inches(1.8), Inches(0.3),
        f"{conf}", font_size=13, color=GRAY, alignment=PP_ALIGN.RIGHT)

add_shape(slide, Inches(0.8), Inches(6.5), Inches(11.7), Inches(0.8), CARD_BG, RGBColor(0xDD, 0xDD, 0xDD))
add_text_box(slide, Inches(1.1), Inches(6.55), Inches(11.3), Inches(0.6),
    "Result: 9/10 videos correctly classified (90% on this sample). Web app correctly detected uploaded fake video as FAKE with 87.35% confidence.",
    font_size=16, color=DARK, bold=True, alignment=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 13: CONCLUSION
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(5), Inches(0.8),
    "Conclusion & Future Work", font_size=36, color=BLUE, bold=True)

card = add_shape(slide, Inches(0.8), Inches(1.5), Inches(11.7), Inches(2.5), CARD_BG, RGBColor(0xDD, 0xDD, 0xDD))
add_text_box(slide, Inches(1.1), Inches(1.7), Inches(11), Inches(0.5),
    "What We Achieved", font_size=22, color=DARK, bold=True)
add_bullet_text(slide, Inches(1.1), Inches(2.3), Inches(11), Inches(1.5), [
    "End-to-end deepfake detection pipeline: Video -> Face Detection -> Feature Extraction -> Classification",
    "78.48% video-level test accuracy with AUC of 0.8622 using honest training methodology (identical face detector for both classes)",
    "Consistent face extraction across training and inference eliminates detection artifact bias",
    "Flask web application for real-time video upload and classification",
], font_size=16, color=GRAY)

card = add_shape(slide, Inches(0.8), Inches(4.3), Inches(11.7), Inches(2.5), CARD_BG, RGBColor(0xDD, 0xDD, 0xDD))
add_text_box(slide, Inches(1.1), Inches(4.5), Inches(11), Inches(0.5),
    "Future Improvements", font_size=22, color=DARK, bold=True)
add_bullet_text(slide, Inches(1.1), Inches(5.1), Inches(11), Inches(1.5), [
    "Fine-tune ResNet50 end-to-end instead of using as fixed feature extractor",
    "Use temporal features (LSTM/3D CNN) instead of per-frame averaging",
    "Expand dataset with more deepfake variants (DFDC, FaceShifter, etc.)",
    "Add video preprocessing to handle different compression levels (c23, c40)",
    "Ensemble multiple models for improved accuracy",
], font_size=16, color=GRAY)

# ============================================================
# SAVE
# ============================================================
output_path = "/home/syed/BNU_PROJECT/Deepfake_Detector_Presentation.pptx"
prs.save(output_path)
print(f"Presentation saved to {output_path}")
print(f"Total slides: {len(prs.slides)}")
